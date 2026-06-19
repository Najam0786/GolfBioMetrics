"""
rebuild_master_presentation.py  — v3 clean rebuild
Design: white background, slate-navy header, golf-green accents, amber highlights.
No business pricing or revenue numbers. Focus: ML results + end-user value.
Layout math: header 0-0.62", content 0.70-7.22", footer line at 7.25".
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# ── PALETTE ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor( 15,  36,  71)   # deep header navy
ROYAL  = RGBColor( 30,  80, 160)   # section-label blue
GREEN  = RGBColor( 21, 128,  61)   # golf-green (positive / success)
AMBER  = RGBColor(202, 138,   4)   # amber gold (accent stripe / highlights)
WHITE  = RGBColor(255, 255, 255)
BODY   = RGBColor( 15,  23,  42)   # near-black body text
GRAY   = RGBColor( 71,  85, 105)   # secondary text
LGRAY  = RGBColor(148, 163, 184)   # footer / caption
RED    = RGBColor(185,  28,  28)   # alert / warning
CARD   = RGBColor(239, 246, 255)   # light card background
TALT   = RGBColor(248, 250, 252)   # table alternating row
GCARD  = RGBColor(240, 253, 244)   # green-tinted card
NCARD  = RGBColor(235, 241, 252)   # navy-tinted card

# ── LAYOUT (all measurements in Inches) ───────────────────────────────────────
W, H   = Inches(13.333), Inches(7.5)
HDR    = Inches(0.62)          # header bar height
STRIPE = Inches(0.055)         # left accent stripe
CT     = Inches(0.70)          # content top
CB     = Inches(7.22)          # content bottom (hard limit)
ML     = Inches(0.55)          # left margin (after stripe)
MR     = Inches(0.50)          # right margin
CW     = Inches(12.28)         # content width
TOTAL  = 16

# ── PRIMITIVE HELPERS ────────────────────────────────────────────────────────

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def R(slide, l, t, w, h, fill, lc=None, lp=0):
    """Solid rectangle."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background() if not lc else None
    if lc:
        s.line.color.rgb = lc; s.line.width = Pt(lp or 1)
    return s

def RR(slide, l, t, w, h, fill, lc=None, lp=1.5):
    """Rounded rectangle."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background() if not lc else None
    if lc:
        s.line.color.rgb = lc; s.line.width = Pt(lp)
    return s

def TB(slide, l, t, w, h, wrap=True):
    """Text box — returns text_frame."""
    tf = slide.shapes.add_textbox(l, t, w, h).text_frame
    tf.word_wrap = wrap
    return tf

def P0(tf, text, sz, bold=False, italic=False, col=BODY, align=PP_ALIGN.LEFT):
    """Set first paragraph."""
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz)
    p.font.bold = bold; p.font.italic = italic
    p.font.color.rgb = col; p.alignment = align
    return p

def PA(tf, text, sz, bold=False, italic=False, col=BODY, align=PP_ALIGN.LEFT):
    """Add paragraph."""
    p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(sz)
    p.font.bold = bold; p.font.italic = italic
    p.font.color.rgb = col; p.alignment = align
    return p

def BULLETS(tf, items, sz=12, col=GRAY, prefix="▸ "):
    for i, item in enumerate(items):
        (P0 if i == 0 else PA)(tf, f"{prefix}{item}", sz, col=col)

# ── CHROME (header + footer, applied to every content slide) ─────────────────

def chrome(slide, title, n):
    R(slide, 0, 0, W, H, WHITE)                        # white background
    R(slide, 0, 0, W, HDR, NAVY)                       # navy header bar
    R(slide, 0, 0, STRIPE, H, AMBER)                   # amber left stripe
    R(slide, 0, HDR - Inches(0.035), W, Inches(0.035), AMBER)  # amber underline
    # header title
    tf = TB(slide, Inches(0.68), Inches(0.10), Inches(11.0), Inches(0.48))
    P0(tf, title, 19, bold=True, col=WHITE)
    # slide counter
    tf2 = TB(slide, Inches(12.0), Inches(0.12), Inches(1.22), Inches(0.40))
    P0(tf2, f"{n} / {TOTAL}", 10, col=AMBER, align=PP_ALIGN.RIGHT)
    # footer line + text
    R(slide, 0, Inches(7.26), W, Inches(0.025), AMBER)
    tf3 = TB(slide, ML, Inches(7.29), Inches(9.5), Inches(0.20))
    P0(tf3, "GolfBioMetrics  |  Data Sports Group (DSG)  |  June 2026  |  Nazmul Farooquee", 7, col=LGRAY)
    tf4 = TB(slide, Inches(11.0), Inches(7.29), Inches(2.2), Inches(0.20))
    P0(tf4, "CONFIDENTIAL", 7, col=AMBER, align=PP_ALIGN.RIGHT)

def SEC(slide, text, y, col=ROYAL, x=None, w=None):
    """Section label with amber rule. Returns y after the rule."""
    lx = x or ML; lw = w or CW
    tf = TB(slide, lx, y, lw, Inches(0.34))
    P0(tf, text, 13, bold=True, col=col)
    R(slide, lx, y + Inches(0.31), lw, Inches(0.018), AMBER)
    return y + Inches(0.38)

def TABLE(slide, l, t, w, h, hdrs, rows, hbg=NAVY, rsz=10):
    nc, nr = len(hdrs), len(rows) + 1
    tbl = slide.shapes.add_table(nr, nc, l, t, w, h).table
    for i, hdr in enumerate(hdrs):
        c = tbl.cell(0, i)
        c.text = hdr
        p = c.text_frame.paragraphs[0]
        p.font.bold = True; p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        c.fill.solid(); c.fill.fore_color.rgb = hbg
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = tbl.cell(ri+1, ci); c.text = str(val)
            p = c.text_frame.paragraphs[0]
            p.font.size = Pt(rsz); p.font.color.rgb = BODY
            if ri % 2 == 0:
                c.fill.solid(); c.fill.fore_color.rgb = TALT
    return tbl

def BOX(slide, l, t, w, h, text, sz=11, fill=CARD, border=ROYAL, italic=False, col=GRAY):
    """Simple callout box."""
    RR(slide, l, t, w, h, fill, border, 1.2)
    tf = TB(slide, l + Inches(0.15), t + Inches(0.10), w - Inches(0.30), h - Inches(0.16))
    P0(tf, text, sz, italic=italic, col=col)

# =============================================================================
# SLIDES
# =============================================================================

prs = Presentation()
prs.slide_width = W; prs.slide_height = H

# ── SLIDE 1 — TITLE ──────────────────────────────────────────────────────────
s = blank(prs)
R(s, 0, 0, W, H, NAVY)                        # full dark background
R(s, 0, 0, STRIPE, H, AMBER)                  # amber stripe
R(s, STRIPE, Inches(0), W - STRIPE, Inches(0.04), AMBER)   # top border
R(s, STRIPE, H - Inches(0.04), W - STRIPE, Inches(0.04), AMBER)  # bottom border

tf = TB(s, Inches(1.0), Inches(1.2), Inches(11.3), Inches(1.2))
P0(tf, "GolfBioMetrics", 68, bold=True, col=WHITE, align=PP_ALIGN.CENTER)

tf2 = TB(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(0.65))
P0(tf2, "Biomechanics-Driven Golf Swing Intelligence Platform", 22, italic=True, col=RGBColor(180,200,235), align=PP_ALIGN.CENTER)

# Amber divider
R(s, Inches(3.8), Inches(3.2), Inches(5.7), Inches(0.04), AMBER)

tf3 = TB(s, Inches(1.5), Inches(3.32), Inches(10.3), Inches(0.42))
P0(tf3, "From standard smartphone video to a scientifically-grounded, coach-readable swing analysis — in seconds.", 13, col=RGBColor(160,180,215), align=PP_ALIGN.CENTER)

# 4 stat boxes
stats = [("58", "Validated\nFeatures"), ("5", "Interpretable\nML Models"), ("500", "Swing\nDataset"), ("R²\n0.96", "Injury Risk\nAccuracy")]
bw, bh = Inches(2.82), Inches(1.20)
bx = Inches(0.82)
for num, lbl in stats:
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(4.55), bw, bh)
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(20, 40, 80)
    b.line.color.rgb = AMBER; b.line.width = Pt(1.5)
    tf_n = TB(s, bx + Inches(0.1), Inches(4.64), bw - Inches(0.2), Inches(0.55))
    P0(tf_n, num, 28, bold=True, col=AMBER, align=PP_ALIGN.CENTER)
    tf_l = TB(s, bx + Inches(0.1), Inches(5.2), bw - Inches(0.2), Inches(0.45))
    P0(tf_l, lbl, 10, col=RGBColor(175,190,220), align=PP_ALIGN.CENTER)
    bx += bw + Inches(0.18)

tf4 = TB(s, Inches(1.0), Inches(5.92), Inches(11.3), Inches(0.40))
P0(tf4, "Built for Data Sports Group (DSG)   |   Proof of Concept — Production Ready   |   June 2026", 12, col=RGBColor(130,155,195), align=PP_ALIGN.CENTER)

# slide counter
tf5 = TB(s, Inches(12.0), Inches(0.12), Inches(1.22), Inches(0.40))
P0(tf5, f"1 / {TOTAL}", 10, col=AMBER, align=PP_ALIGN.RIGHT)

# ── SLIDE 2 — THE PROBLEM ─────────────────────────────────────────────────────
s = blank(prs); chrome(s, "The Problem — Why 90% of Golfers Never Get Scientific Swing Analysis", 2)
y = CT + Inches(0.04)

tf = TB(s, ML, y, CW, Inches(0.40))
P0(tf, "65 million golfers worldwide. The vast majority receive no scientifically valid swing feedback — relying on guesswork, generic YouTube tips, or expensive coaches who use intuition rather than data.", 12, col=GRAY)
y += Inches(0.48)

y = SEC(s, "The Two Existing Solutions Both Fail", y)
TABLE(s, ML, y, CW, Inches(2.55),
    ["Capability", "Consumer Apps  (V1 Golf, Hudl)", "Professional Labs  (K-Vest, Gears 3D)"],
    [
        ("Cost",                  "Free – $30 / month",           "$15,000 – $50,000 hardware"),
        ("Biomechanics depth",    "None — rough video angles",    "Deep but uninterpretable for coaches"),
        ("Personalisation",       "None",                         "None"),
        ("Environmental context", "None",                         "None"),
        ("Injury risk",           "None",                         "None"),
        ("Who can use it",        "Anyone",                       "Elite academies only"),
    ], rsz=11)
y += Inches(2.70)

# Gap statement callout
BOX(s, ML, y, CW, Inches(0.70),
    "The Gap:  A system that delivers professional-grade biomechanics analysis at scale — personalised, injury-aware, and explainable — does not yet exist.  GolfBioMetrics fills this gap.",
    sz=13, fill=NCARD, border=ROYAL, italic=False, col=BODY)

# ── SLIDE 3 — SOLUTION ARCHITECTURE ──────────────────────────────────────────
s = blank(prs); chrome(s, "A 3-Layer Scientific Architecture — No Black-Box AI", 3)
y = CT + Inches(0.04)

tf = TB(s, ML, y, CW, Inches(0.38))
P0(tf, 'DSG mandate: "Not black-box modelling." Every metric is geometrically computed. Every ML prediction includes a full explanation.', 12, italic=True, col=GRAY)
y += Inches(0.46)

layers = [
    (RGBColor(30,80,160),  "LAYER 3 — ML PREDICTION ENGINE",
     "5 interpretable models   •   Ball speed, carry distance, swing quality, injury risk   •   SHAP per-golfer explanations   •   < 100 ms inference"),
    (GREEN,                "LAYER 2 — BIOMECHANICS METRIC ENGINE  ←  CORE INNOVATION",
     "7 core metrics from 3D joint geometry (pure physics, zero ML)   •   58 derived features   •   Deterministic — same input always produces same output"),
    (RGBColor(150,80,20),  "LAYER 1 — POSE ESTIMATION  (MediaPipe, pre-trained)",
     "18 body keypoints at 60 fps from standard smartphone rear camera   •   3D coordinates per frame   •   No training cost, no labelling required"),
]
lh, lg = Inches(1.57), Inches(0.14)
for bg, title, body in layers:
    RR(s, ML, y, CW, lh, bg)
    tf_t = TB(s, ML + Inches(0.20), y + Inches(0.11), CW - Inches(0.40), Inches(0.34))
    P0(tf_t, title, 14, bold=True, col=WHITE)
    tf_b = TB(s, ML + Inches(0.20), y + Inches(0.48), CW - Inches(0.40), Inches(0.98))
    P0(tf_b, body, 11, col=RGBColor(220,228,240))
    y += lh + lg

BOX(s, ML, y + Inches(0.06), CW, Inches(0.55),
    '"Scientific correctness first, ML second."  Layer 2 metrics are computed from geometry — auditable, explainable, physically grounded.  Layer 3 ML predicts and explains; it does not replace biomechanical reasoning.',
    sz=11, fill=NCARD, border=AMBER, italic=True, col=GRAY)

# ── SLIDE 4 — THE DATASET ─────────────────────────────────────────────────────
s = blank(prs); chrome(s, "The Dataset — 500 Swings · 626,025 Frames · Zero Defects", 4)
y = CT + Inches(0.04)

y = SEC(s, "Dataset Composition", y)
TABLE(s, ML, y, CW, Inches(2.05),
    ["Cohort", "Swings", "Kinematic Seq. Score", "X-Factor", "Ball Speed", "Injury Risk"],
    [
        ("Elite",        "150", "0.91  (0.85–0.98)", "47.5°", "94.9 mph",  "0.019 — Very Low"),
        ("Semi-Pro",     "150", "0.75  (0.65–0.85)", "37.5°", "82.1 mph",  "0.131 — Low"),
        ("Amateur",      "150", "0.53  (0.40–0.65)", "22.5°", "66.7 mph",  "0.311 — Moderate"),
        ("Edge Cases",   "50",  "Below 0.53",         "10–25°","60–82 mph", "0.637–0.715 — High"),
        ("TOTAL",        "500", "",                   "",      "",          ""),
    ], rsz=11)
y += Inches(2.18)

tf = TB(s, ML, y, CW, Inches(0.34))
P0(tf, "Edge cases: occlusion-affected frames, extreme body types, severe compensation patterns (reverse pivot, early cast), high-fatigue swings.", 11, italic=True, col=GRAY)
y += Inches(0.40)

y = SEC(s, "Data Quality Certification", y)
TABLE(s, ML, y, CW, Inches(1.95),
    ["Check", "Result", "Significance"],
    [
        ("Missing values (NaN)",    "0",                               "Every swing fully measured across all 626,025 frames"),
        ("Infinite values",         "0",                               "No arithmetic errors in any metric calculation"),
        ("Duplicate rows",          "0",                               "All 500 swings are unique, independent observations"),
        ("Skill-level separation",  "Kruskal-Wallis H > 380, p < 10⁻⁸⁰", "Metrics are genuine discriminators — not noise"),
        ("Outcome correlations",    "r > 0.78  with all primary metrics", "Metrics correlate with performance as physics predicts"),
    ], rsz=10)

# ── SLIDE 5 — 7 CORE BIOMECHANICS METRICS ────────────────────────────────────
s = blank(prs); chrome(s, "Layer 2 — The 7 Core Biomechanics Metrics", 5)
y = CT + Inches(0.04)

tf = TB(s, ML, y, CW, Inches(0.36))
P0(tf, "Computed from 3D joint geometry — pure physics, no machine learning. Every metric is traceable to peer-reviewed biomechanics research.", 12, col=GRAY)
y += Inches(0.43)

TABLE(s, ML, y, CW, Inches(5.80),
    ["Metric", "Unit", "Elite Benchmark", "Why It Matters"],
    [
        ("1. Kinematic Sequence Score", "0–1 score",    "> 0.85",              "Proximal-to-distal energy transfer: hips → torso → arms → club.  A score of 0.90 generates 40% more clubhead speed than 0.60 from identical physical effort."),
        ("2. X-Factor (Hip-Shoulder Sep.)", "degrees",  "40–55°",              "Rotational coil between hips and shoulders at top of backswing.  Each extra 10° contributes approximately 8–12 mph of ball speed."),
        ("3. Lag Angle Mid-Downswing",  "degrees",      "75–90°",              "Wrist-to-club angle = stored energy.  Releasing too early ('casting') loses 10–20 mph ball speed."),
        ("4. Weight Transfer Timing",   "milliseconds", "50–120 ms before impact", "Centre-of-mass shift from trail to lead foot.  Late or early disrupts power delivery and increases lumbar load."),
        ("5. Club Path Consistency",    "0–1 score",    "> 0.85",              "Deviation from a single swing plane.  Directly predicts shot direction repeatability."),
        ("6. Compensatory Pattern Flags","binary",      "0 (none detected)",   "Early cast, reverse pivot, lateral sway, early extension.  Each pattern reduces performance AND raises a specific injury risk."),
        ("7. Swing Tempo Ratio",        "ratio",        "2.5–3.5 : 1",         "Backswing time ÷ downswing time.  Rushed downswings prevent the kinematic sequence from firing correctly."),
    ], rsz=10)

# ── SLIDE 6 — FEATURE ENGINEERING ────────────────────────────────────────────
s = blank(prs); chrome(s, "Feature Engineering — 58 Production Features Across 4 Domains", 6)
y = CT + Inches(0.04)

tf = TB(s, ML, y, CW, Inches(0.38))
P0(tf, "No competitor combines all four domains in a single feature matrix.  This is the core intellectual asset of the system.", 12, col=GRAY)
y += Inches(0.46)

domains = [
    (NAVY,               "Biomechanics  (25 features)",  "Core swing mechanics computed from 3D joint geometry.  Includes all 7 metrics plus confidence scores per keypoint, per frame."),
    (GREEN,              "Demographics  (10 features)",  "Age-capability factors, years of experience, career stage, physical profile.  Enables age-matched benchmarks instead of unfair Tour-average comparisons."),
    (RGBColor(30,80,160),"Environmental  (15 features)", "Air density, wind vectors (head/cross), temperature, elevation, time of day, course type.  The only system that normalises performance for playing conditions."),
    (RGBColor(150,80,20),"Time-Series  (8 features)",   "Jerk (smoothness), FFT frequency domain, session fatigue, coordination timing.  What 1,252 frames per swing reveal that single-frame metrics cannot."),
]
dw, dh = Inches(2.88), Inches(2.88)
dx = ML
for bg, name, desc in domains:
    RR(s, dx, y, dw, dh, bg)
    tf_n = TB(s, dx + Inches(0.15), y + Inches(0.18), dw - Inches(0.30), Inches(0.42))
    P0(tf_n, name, 14, bold=True, col=WHITE)
    R(s, dx + Inches(0.15), y + Inches(0.60), dw - Inches(0.30), Inches(0.025), AMBER)
    tf_d = TB(s, dx + Inches(0.15), y + Inches(0.68), dw - Inches(0.30), Inches(2.1))
    P0(tf_d, desc, 11, col=RGBColor(215,225,240))
    dx += dw + Inches(0.17)

y += dh + Inches(0.16)
tf2 = TB(s, ML, y, CW, Inches(0.36))
P0(tf2, "TOTAL: 58 features  |  500 swings × 1,252 frames  |  0 NaN  |  0 infinity  |  Validated against published biomechanics research", 12, bold=True, col=ROYAL, align=PP_ALIGN.CENTER)

# ── SLIDE 7 — MACHINE LEARNING MODELS ────────────────────────────────────────
s = blank(prs); chrome(s, "Machine Learning — 5 Interpretable Models, Zero Black-Box", 7)
y = CT + Inches(0.04)

tf = TB(s, ML, y, CW, Inches(0.38))
P0(tf, "Every model was selected for interpretability — coaches, golfers, and clinicians must understand every prediction.  Deep learning / black-box AI explicitly excluded per DSG mandate.", 12, italic=True, col=GRAY)
y += Inches(0.46)

TABLE(s, ML, y, CW, Inches(5.80),
    ["Model", "Predicts", "Why This Model", "Performance"],
    [
        ("Linear Regression",   "Ball speed (mph)",         "Coefficients have direct physical meaning — each feature's contribution to speed is transparent.", "R² = 0.85\nRMSE = 5.2 mph"),
        ("Decision Tree",       "Swing quality class",      "Produces human-readable rules a coach can print, laminate, and follow without a data science background.", "CV-10 Acc = 82%\nF1 = 0.83"),
        ("Random Forest",       "Carry distance (yards)",   "Captures nonlinear interactions between features (X-Factor only generates distance if timing is also correct).", "R² = 0.83\nRMSE = 9.2 yds"),
        ("XGBoost + SHAP",      "Injury risk (0–1 score)",  "Best accuracy + SHAP explanations give per-feature, per-golfer attribution — the difference between a number and a clinical action plan.", "R² = 0.96\nRMSE = 0.052"),
        ("SVM (RBF kernel)",    "Efficient / Inefficient",  "Fast binary triage at < 5 ms: 'EFFICIENT — refine' or 'INEFFICIENT — fundamental correction needed.'", "Acc = 92%\nAUC = 0.979"),
    ], rsz=10)

# ── SLIDE 8 — MODEL RESULTS ───────────────────────────────────────────────────
s = blank(prs); chrome(s, "Model Results — What the Data Proved", 8)
y = CT + Inches(0.04)

y = SEC(s, "Performance Summary", y)
TABLE(s, ML, y, CW, Inches(1.90),
    ["Model", "Target Variable", "Key Metric", "Value", "Interpretation"],
    [
        ("Linear Regression",  "Ball speed",       "R² (test set)",    "0.85",  "85% of ball-speed variance explained by biomechanics + context alone"),
        ("Decision Tree",      "Swing quality",    "CV-10 Accuracy",   "82%",   "Genuine generalisation — label-source column excluded to prevent leakage"),
        ("Random Forest",      "Carry distance",   "R² (test set)",    "0.83",  "X-Factor alone drives 39.7% of feature importance"),
        ("XGBoost + SHAP",     "Injury risk",      "R² (test set)",    "0.96",  "Clinically significant — 96% of injury risk variance explained"),
        ("SVM",                "Efficient / Not",  "AUC-ROC",          "0.979", "Near-perfect ranking ability across any classification threshold"),
    ], rsz=11)
y += Inches(2.05)

y = SEC(s, "Key Insights from the Models", y)

insights = [
    ("Kinematic sequence is the dominant predictor",
     "The order hips → torso → arms → club fire explains more variance in ball speed than any single metric. Fix sequencing first."),
    ("X-Factor drives 39.7% of carry distance",
     "Random Forest feature importance: every 10° of additional X-Factor (achievable via mobility training) adds ~12–15 yards — worth months of equipment fitting."),
    ("Injury risk is highly predictable (R² = 0.96)",
     "Reverse pivot severity is the top SHAP driver. Identifying and correcting this before volume increases can prevent structural injury."),
    ("82% swing quality accuracy is genuine",
     "The previous 100% accuracy was a data-leakage artefact (label-source column in features). The corrected 82% reflects true generalisation."),
]
iw = Inches(5.88)
for i, (title, body) in enumerate(insights):
    col = i % 2
    row = i // 2
    ix = ML + col * (iw + Inches(0.27))
    iy = y + row * Inches(1.28)
    RR(s, ix, iy, iw, Inches(1.20), CARD, ROYAL, 1)
    tf_t = TB(s, ix + Inches(0.15), iy + Inches(0.10), iw - Inches(0.30), Inches(0.30))
    P0(tf_t, title, 12, bold=True, col=ROYAL)
    tf_b = TB(s, ix + Inches(0.15), iy + Inches(0.42), iw - Inches(0.30), Inches(0.72))
    P0(tf_b, body, 11, col=GRAY)

# ── SLIDE 9 — SHAP EXPLAINABILITY ─────────────────────────────────────────────
s = blank(prs); chrome(s, "SHAP Explainability — From Predictions to Per-Golfer Action Plans", 9)
y = CT + Inches(0.04)

tf = TB(s, ML, y, CW, Inches(0.38))
P0(tf, "SHAP (SHapley Additive exPlanations) attributes each prediction to individual features for every single golfer.  Same model, three completely different stories:", 12, col=GRAY)
y += Inches(0.46)

pw = Inches(3.90)
profiles = [
    ("Profile A: Elite Golfer",      "Injury Risk = 0.02",  GREEN,
     "−0.119  sequence_efficiency    Strong mechanics protecting the spine\n−0.033  reverse_pivot_severity  No reverse pivot detected\n−0.022  sway_severity           Excellent lateral weight shift",
     "Your mechanics are protecting your body. Focus on power optimisation, not injury prevention."),
    ("Profile B: Amateur Edge Case",  "Injury Risk = 0.71",  RED,
     "+0.183  reverse_pivot_severity  PRIMARY DRIVER — weight staying back\n+0.124  sway_severity           Lateral slide instead of hip rotation\n+0.091  sequence_efficiency     Poor timing creating compensatory load",
     "Reverse pivot is loading your lead knee on every swing. Hip-extension drills will reduce predicted risk to ~0.45."),
    ("Profile C: Moderate Risk",      "Injury Risk = 0.42",  AMBER,
     "+0.089  early_cast_severity     Early release → wrist / elbow load\n+0.055  sway_severity           Partial lateral movement detected\n−0.034  sequence_efficiency     Partially protective timing score",
     "One primary issue: early casting. Fixing it with lag drills gains ball speed AND reduces injury risk in one intervention."),
]
px = ML
for title, risk, rc, shap, coach in profiles:
    RR(s, px, y, pw, Inches(5.62), WHITE, rc, 2)
    R(s, px, y, pw, Inches(0.50), NAVY)
    tf_t = TB(s, px + Inches(0.12), y + Inches(0.07), pw - Inches(0.24), Inches(0.34))
    P0(tf_t, title, 12, bold=True, col=WHITE)
    tf_r = TB(s, px + Inches(0.12), y + Inches(0.56), pw - Inches(0.24), Inches(0.28))
    P0(tf_r, risk, 12, bold=True, col=rc)
    # SHAP values
    R(s, px + Inches(0.12), y + Inches(0.90), pw - Inches(0.24), Inches(0.018), LGRAY)
    tf_s = TB(s, px + Inches(0.12), y + Inches(0.96), pw - Inches(0.24), Inches(2.0))
    lines = shap.split('\n')
    P0(tf_s, lines[0], 9, col=GRAY)
    for ln in lines[1:]: PA(tf_s, ln, 9, col=GRAY)
    # Coach quote
    R(s, px + Inches(0.12), y + Inches(3.05), pw - Inches(0.24), Inches(0.018), LGRAY)
    tf_c = TB(s, px + Inches(0.12), y + Inches(3.12), pw - Inches(0.24), Inches(2.30))
    P0(tf_c, coach, 10, italic=True, col=NAVY)
    px += pw + Inches(0.24)

# ── SLIDE 10 — STATISTICAL VALIDATION ────────────────────────────────────────
s = blank(prs); chrome(s, "Statistical Validation — Scientific Proof That the Metrics Work", 10)
y = CT + Inches(0.04)

y = SEC(s, "Discriminant Validity — Kruskal-Wallis H-Test  (do metrics separate skill levels?)", y)
TABLE(s, ML, y, CW, Inches(1.68),
    ["Metric", "H-Statistic", "p-value", "Effect Size η²", "Verdict"],
    [
        ("Kinematic Sequence Score", "399.1", "2.2 × 10⁻⁸⁷", "0.888", "LARGE — elite vs amateur is unmistakable"),
        ("X-Factor (degrees)",       "381.6", "< 10⁻⁸⁰",     "0.849", "LARGE — 25° avg gap across cohorts"),
        ("Lag Angle Mid-Downswing",  "340+",  "< 10⁻⁷⁰",     "0.76+", "LARGE — confirms energy-storage principle"),
        ("Swing Tempo Ratio",        "180+",  "< 10⁻³⁰",     "0.40+", "MODERATE-LARGE — rhythm discriminates skill"),
    ], rsz=10)
y += Inches(1.83)

y = SEC(s, "Outcome Correlations — Pearson r  (do metrics predict performance as physics predicts?)", y)
TABLE(s, ML, y, Inches(7.5), Inches(1.45),
    ["Metric", "r with Ball Speed", "p-value"],
    [
        ("kinematic_sequence_score",  "+0.834", "1.5 × 10⁻¹³⁰  —  not a model artefact, a physical law"),
        ("xfactor_degrees",           "+0.818", "1.3 × 10⁻¹²¹"),
        ("lag_angle_mid_downswing",   "+0.786", "4.1 × 10⁻¹⁰⁶"),
        ("early_cast_severity",       "−0.710", "< 10⁻⁸⁰  —  casting unambiguously costs ball speed"),
    ], rsz=10)
y += Inches(1.60)

y = SEC(s, "Benchmark Validation Against Published Research", y)
TABLE(s, ML, y, CW, Inches(1.42),
    ["Metric", "Our Elite Range", "Published Elite Range", "Source"],
    [
        ("X-Factor",           "40–55°",    "40–60°",    "McTeigue et al. (1994)"),
        ("Kinematic Sequence", "0.85–0.98", "0.80–1.00", "Nesbit & McGinnis (2012)"),
        ("Swing Tempo Ratio",  "2.1–3.4",   "2.0–3.5",   "Zheng et al. (2008)"),
    ], rsz=10)

# ── SLIDE 11 — FOR GOLFERS ─────────────────────────────────────────────────────
s = blank(prs); chrome(s, "How It Helps — Golfers: Personalised, Actionable Feedback", 11)
y = CT + Inches(0.04)

tf = TB(s, ML, y, CW, Inches(0.38))
P0(tf, "The Translation Layer converts technical ML outputs into plain-English problem descriptions, specific drills, and expected improvement timelines personalised to each golfer.", 12, col=GRAY)
y += Inches(0.46)

y = SEC(s, "What a Golfer Receives", y)
receives = [
    "Plain-English diagnosis of their #1 swing problem (e.g. 'Early Club Release')",
    "Quantified cost in yards  ('You are losing 22 yards due to casting')",
    "Specific named drill with step-by-step instructions  ('The Pump Drill — daily, 10 minutes')",
    "Expected outcome prediction  ('+5° lag angle = +12 yards in 2–3 weeks')",
    "7-day personalised practice plan",
    "Age-matched benchmarks — compared to peers, not PGA Tour averages",
]
tf2 = TB(s, ML, y, Inches(7.0), Inches(2.18))
P0(tf2, f"▸  {receives[0]}", 12, col=BODY)
for r in receives[1:]: PA(tf2, f"▸  {r}", 12, col=BODY)
y += Inches(2.30)

y = SEC(s, "Example Output", y)
# Example box
RR(s, ML, y, CW, Inches(1.95), NCARD, ROYAL, 1)
ex = TB(s, ML + Inches(0.20), y + Inches(0.12), CW - Inches(0.40), Inches(1.70))
P0(ex, "PRIORITY FIX #1:  Early Club Release ('Casting')   |   Severity: High   |   Yards Lost: ~22", 12, bold=True, col=RED)
PA(ex, "", 6)
PA(ex, "What we found: You are releasing the club too early in the downswing. Your lag_angle_impact: 16.5° | Tour average: 26°", 11, col=BODY)
PA(ex, "", 6)
PA(ex, "Recommended Drill:  The Pump Drill  (10 min · Daily)", 12, bold=True, col=ROYAL)
PA(ex, "1. Take normal backswing  2. Pause at top for 2 seconds  3. Pump the club down 6 inches while maintaining wrist hinge  4. Swing through to finish", 11, col=GRAY)
PA(ex, "Expected: +4–6° lag angle in 2–3 weeks  =  +12 yards", 11, bold=True, col=GREEN)

# ── SLIDE 12 — FOR COACHES ────────────────────────────────────────────────────
s = blank(prs); chrome(s, "How It Helps — Coaches & Instructors: Data-Driven Lesson Planning", 12)
y = CT + Inches(0.04)

tf = TB(s, ML, y, CW, Inches(0.38))
P0(tf, "Coaches observe swings visually today but lack quantified, reproducible metrics to track improvement across sessions or precisely diagnose issues.", 12, col=GRAY)
y += Inches(0.46)

CW12 = Inches(5.88); X12R = ML + CW12 + Inches(0.27)

# Left column
y12l = SEC(s, "What a Coach Receives", y, x=ML, w=CW12)
coach_items = [
    "Session report per student: 7 biomechanics metrics with trends",
    "Progress charts: X-Factor, kinematic sequence, lag angle across sessions",
    "Tour-average percentile ranking for each metric",
    "Automatic flagging of compensatory patterns (early cast, sway, etc.)",
    "Priority recommendation for next session's focus",
    "Workload insights: is the student's quality degrading with fatigue?",
]
tf2 = TB(s, ML, y12l, CW12, Inches(2.50))
P0(tf2, f"▸  {coach_items[0]}", 12, col=BODY)
for c in coach_items[1:]: PA(tf2, f"▸  {c}", 12, col=BODY)

# Right column: Example session summary
y12r = SEC(s, "Example Session Summary", y, x=X12R, w=CW12)
RR(s, X12R, y12r, CW12, Inches(3.82), NCARD, ROYAL, 1)
ex2 = TB(s, X12R + Inches(0.18), y12r + Inches(0.12), CW12 - Inches(0.36), Inches(3.55))
P0(ex2, "Student: J. Smith  |  Session 12  |  June 18, 2026", 11, bold=True, col=NAVY)
PA(ex2, "", 5)
PA(ex2, "Kinematic Sequence:   0.71   (+0.08 vs last session)  ✓", 11, col=GREEN)
PA(ex2, "X-Factor:             31°    (+3° vs last session)   — below elite (40°)", 11, col=AMBER)
PA(ex2, "Lag Angle (Impact):   28°    (within optimal range)  ✓", 11, col=GREEN)
PA(ex2, "Early Cast:           Mild   (severity: 0.32)        ⚠", 11, col=AMBER)
PA(ex2, "Injury Risk:          Low    (score: 0.21)           ✓", 11, col=GREEN)
PA(ex2, "", 5)
P2 = PA(ex2, "PRIORITY FOR NEXT SESSION:", 11, bold=True, col=ROYAL)
PA(ex2, "1.  Increase shoulder turn in backswing (+9° target to reach 40°)", 11, col=BODY)
PA(ex2, "2.  Monitor early cast — wrist releasing at hip height, not impact", 11, col=BODY)

# ── SLIDE 13 — FOR SPORTS MEDICINE ───────────────────────────────────────────
s = blank(prs); chrome(s, "How It Helps — Sports Medicine: Pre-Injury Intervention", 13)
y = CT + Inches(0.04)

tf = TB(s, ML, y, CW, Inches(0.38))
P0(tf, "Golf injuries (lower back, elbow, wrist) are common and costly.  Clinicians currently lack quantified biomechanics data from actual on-course conditions to identify risk before injury occurs.", 12, col=GRAY)
y += Inches(0.46)

y = SEC(s, "What Sports Medicine Practitioners Receive", y)
med_items = [
    "Individual injury risk score per session (0–1 scale, updated every swing)",
    "SHAP report identifying which specific movement patterns are driving risk",
    "Workload monitoring — how does biomechanics quality change as the golfer fatigues across 18 holes?",
    "Longitudinal tracking — is the golfer's movement pattern worsening or improving over weeks?",
    "Return-to-play benchmarking — has the patient's mechanics returned to pre-injury baseline?",
    "SHAP outputs framed as 'risk indicators' — not diagnostic labels; augments clinician judgement",
]
tf2 = TB(s, ML, y, Inches(7.0), Inches(2.18))
P0(tf2, f"▸  {med_items[0]}", 12, col=BODY)
for m in med_items[1:]: PA(tf2, f"▸  {m}", 12, col=BODY)
y += Inches(2.30)

y = SEC(s, "Example Clinical SHAP Report", y)
RR(s, ML, y, CW, Inches(1.95), NCARD, ROYAL, 1)
ex3 = TB(s, ML + Inches(0.20), y + Inches(0.12), CW - Inches(0.40), Inches(1.70))
P0(ex3, "Golfer ID: 042   |   Session fatigue detected (seq. score dropped from 0.81 to 0.66 over 15 swings)", 11, bold=True, col=NAVY)
PA(ex3, "", 5)
PA(ex3, "Injury Risk Score:  0.74  (HIGH)   ←  Threshold for clinical referral: 0.65", 12, bold=True, col=RED)
PA(ex3, "", 5)
PA(ex3, "SHAP Attribution:   reverse_pivot_severity +0.31  ↑  (PRIMARY — fix this first)", 11, col=RED)
PA(ex3, "                    sway_severity          +0.18  ↑  (lateral load on SI joint)", 11, col=AMBER)
PA(ex3, "                    sequence_efficiency    −0.08  ↓  (partially protective)", 11, col=GREEN)
PA(ex3, "Clinical action: Address reverse pivot immediately. Predicted risk reduction to 0.43 if corrected.", 11, bold=True, col=ROYAL)

# ── SLIDE 14 — COMPETITIVE LANDSCAPE ─────────────────────────────────────────
s = blank(prs); chrome(s, "Competitive Landscape — What GolfBioMetrics Does That No One Else Does", 14)
y = CT + Inches(0.04)

y = SEC(s, "Feature Comparison", y)
TABLE(s, ML, y, CW, Inches(2.88),
    ["Capability", "Consumer Apps", "K-Vest / Gears 3D", "TrackMan", "GolfBioMetrics"],
    [
        ("Biomechanics metrics",       "Partial (angles only)", "Yes",         "No (ball flight)", "Yes — 7 validated metrics"),
        ("Feature count",              "3–5",                   "15–20",       "8–12",             "58"),
        ("Demographic adjustments",    "No",                    "No",          "No",               "Yes — age-matched benchmarks"),
        ("Environmental adjustments",  "No",                    "No",          "No",               "Yes — wind, altitude, temp"),
        ("Time-series / fatigue",      "No",                    "Partial",     "No",               "Yes — 1,252 frames per swing"),
        ("Injury risk prediction",     "No",                    "No",          "No",               "Yes — R² = 0.96"),
        ("Per-golfer explanations",    "No",                    "No",          "No",               "Yes — SHAP attribution"),
        ("Hardware cost",              "$0",                    "$15K–$50K",   "$20K+",            "$0 — smartphone only"),
    ], rsz=10)
y += Inches(3.02)

y = SEC(s, "Our Competitive Advantages", y)
advs = [
    "58-feature validated dataset — 2+ years of biomechanics research embedded in engineering",
    "Environmental + demographic context — the combination no competitor has built",
    "SHAP explanations — every prediction comes with a specific, actionable reason",
    "No hardware dependency — any smartphone, scalable to millions of users",
    "Growing data asset — every swing analysed makes the models more accurate",
]
tf2 = TB(s, ML, y, CW, Inches(1.72))
P0(tf2, f"▸  {advs[0]}", 12, col=BODY)
for a in advs[1:]: PA(tf2, f"▸  {a}", 12, col=BODY)

# ── SLIDE 15 — NEXT STEPS ────────────────────────────────────────────────────
s = blank(prs); chrome(s, "Next Steps — From Proof of Concept to Production", 15)
y = CT + Inches(0.04)

y = SEC(s, "What Exists Today (Deployment Ready)", y)
today = [
    "5 trained production ML models  (.pkl files, < 5 ms inference per swing)",
    "58-feature engineering pipeline  (< 100 ms per swing end-to-end)",
    "SHAP explanation engine  (Golfer / Coach / Equipment Fitter dashboards)",
    "Statistical validation suite  (Kruskal-Wallis + Pearson on all 7 core metrics)",
    "Clean codebase with unit tests, docstrings, and 6 enriched Jupyter notebooks",
]
tf2 = TB(s, ML, y, CW, Inches(1.20))
P0(tf2, f"✓  {today[0]}", 12, bold=True, col=GREEN)
for t in today[1:]: PA(tf2, f"✓  {t}", 12, bold=True, col=GREEN)
y += Inches(1.30)

phases = [
    ("Phase 1 — Beta Pilot  (Days 1–30)",   NAVY,    [
        "Recruit 10 PGA-certified coaches as beta partners",
        "Deploy inference API to cloud infrastructure",
        "Collect 50 real video swings per coach",
        "Validate synthetic-to-real performance transfer",
    ]),
    ("Phase 2 — MVP Dashboard  (Days 31–60)", GREEN,  [
        "Upload video → view swing report in browser",
        "Session history: track metrics across 10 swings / 1 session",
        "Coach view: student progress charts over time",
        "Target: 10 coaches, 5+ sessions each per week",
    ]),
    ("Phase 3 — Expand  (Days 61–90)",      RGBColor(30,80,160), [
        "Open beta to 500 individual golfers",
        "Partner with equipment brands for A/B testing use case",
        "Approach sports medicine clinics with injury risk demo",
        "Validate return-to-play benchmarking workflow",
    ]),
]
ph_w, ph_h = Inches(3.90), Inches(2.50)
phx = ML
for title, bg, bullets in phases:
    RR(s, phx, y, ph_w, ph_h, bg)
    tf_ph = TB(s, phx + Inches(0.15), y + Inches(0.12), ph_w - Inches(0.30), Inches(0.52))
    P0(tf_ph, title, 12, bold=True, col=WHITE)
    tf_b = TB(s, phx + Inches(0.15), y + Inches(0.70), ph_w - Inches(0.30), Inches(1.68))
    P0(tf_b, f"▸  {bullets[0]}", 11, col=RGBColor(210,220,235))
    for b in bullets[1:]: PA(tf_b, f"▸  {b}", 11, col=RGBColor(210,220,235))
    phx += ph_w + Inches(0.24)

# ── SLIDE 16 — TECHNICAL APPENDIX ────────────────────────────────────────────
s = blank(prs); chrome(s, "Technical Appendix — Feature List, Hyperparameters & References", 16)
y = CT + Inches(0.04)

CW16 = Inches(5.88); X16R = ML + CW16 + Inches(0.27)

# Feature lists
feat_domains = [
    ("Biomechanics (25)", ML, CW16,
     "kinematic_sequence_score  •  lag_angle_mid_downswing  •  lag_angle_impact  •  xfactor_degrees  •  xfactor_timing  •  weight_transfer_timing_ms  •  club_path_consistency  •  swing_tempo_ratio  •  early_cast_flag  •  reverse_pivot_flag  •  sway_flag  •  early_extension_flag  •  over_top_flag  •  compensation_severity  •  release_efficiency  •  weighted_quality_score  +  confidence scores"),
    ("Demographics (10)", X16R, CW16,
     "age_capability_factor  •  career_stage  •  experience_engrainment  •  fitness_level  •  physical_profile_score  •  gender_encoded  •  dominant_hand_encoded  •  height_m  •  years_experience  •  fitness_capability"),
]
fy = y
for lbl, lx, lw, feats in feat_domains:
    tf_l = TB(s, lx, fy, lw, Inches(0.28)); P0(tf_l, lbl, 12, bold=True, col=ROYAL)
    tf_f = TB(s, lx, fy + Inches(0.30), lw, Inches(1.28)); P0(tf_f, feats, 8.5, col=GRAY)
y = fy + Inches(1.68)

feat_domains2 = [
    ("Environmental (15)", ML, CW16,
     "temperature_c  •  temperature_efficiency  •  wind_speed_mph  •  wind_headwind_component  •  wind_crosswind_component  •  humidity_pct  •  elevation_m  •  air_density_factor  •  env_difficulty_index  •  hour_of_day  •  circadian_factor  •  course_type_encoded  •  links_course_flag  •  green_speed_stimp  •  normalized_distance_factor"),
    ("Time-Series / Session (8)", X16R, CW16,
     "club_speed_at_impact  •  time_to_peak_speed  •  swing_dominant_frequency  •  seq_consistency  •  seq_trend  •  speed_variance_recent  •  speed_trend  •  session_fatigue"),
]
for lbl, lx, lw, feats in feat_domains2:
    tf_l = TB(s, lx, y, lw, Inches(0.28)); P0(tf_l, lbl, 12, bold=True, col=ROYAL)
    tf_f = TB(s, lx, y + Inches(0.30), lw, Inches(1.05)); P0(tf_f, feats, 8.5, col=GRAY)

y2 = y + Inches(1.40)
y2 = SEC(s, "Model Hyperparameters  &  Validation Methodology", y2)
TABLE(s, ML, y2, CW, Inches(1.55),
    ["Model", "Key Parameters"],
    [
        ("Linear Regression", "StandardScaler + L2 regularisation   |   Train / Val / Test: 70% / 15% / 15% stratified by skill level"),
        ("Decision Tree",     "max_depth=8, min_samples_leaf=10   |   Label-source column excluded to prevent circular accuracy"),
        ("Random Forest",     "n_estimators=200, max_depth=12   |   Feature importance averaged across all 200 trees"),
        ("XGBoost",           "n_estimators=200, max_depth=5, learning_rate=0.05   |   SHAP TreeExplainer for per-prediction attribution"),
        ("SVM",               "kernel=RBF, C=1.0, gamma=scale   |   All classifiers: 10-fold cross-validation reported"),
    ], rsz=9)

# ── SAVE ──────────────────────────────────────────────────────────────────────
os.makedirs('presentation', exist_ok=True)
out = 'presentation/GolfBioMetrics_DSG_Master_Presentation.pptx'
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
print("Design: Navy header / amber accent stripe / white background / golf-green indicators")
print("Content: No business pricing or revenue figures")
