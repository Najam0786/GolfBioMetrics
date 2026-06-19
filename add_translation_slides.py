"""
add_translation_slides.py
Add 2 new slides about Translation Layer to the Master Presentation.
Slides 17-18: Translation Layer and Complete Platform
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

print("=" * 70)
print("ADDING TRANSLATION LAYER SLIDES TO MASTER PRESENTATION")
print("=" * 70)
print()

# Load existing Master Presentation
pptx_path = 'presentation/GolfBioMetrics_DSG_Master_Presentation.pptx'

if not os.path.exists(pptx_path):
    print(f"❌ Error: {pptx_path} not found!")
    print("Make sure the Master Presentation exists.")
    exit(1)

prs = Presentation(pptx_path)
current_slide_count = len(prs.slides)

print(f"✅ Loaded: {pptx_path}")
print(f"📊 Current slides: {current_slide_count}")
print()

# Colors (matching the navy/gold theme)
NAVY = RGBColor(21, 101, 192)  # Primary blue
GOLD = RGBColor(197, 165, 114)  # Gold accent
DARK_GRAY = RGBColor(80, 80, 80)
LIGHT_GRAY = RGBColor(120, 120, 120)

# ============================================================
# SLIDE 17: Translation Layer
# ============================================================
print("Creating Slide 17: Translation Layer...")

slide17 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title_box = slide17.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
tf = title_box.text_frame
tf.text = "The Translation Layer: From ML Metrics to Human Understanding"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = NAVY
p.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide17.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.5))
tf = subtitle_box.text_frame
tf.text = "The Critical Bridge: Making AI Actionable for Golfers, Coaches, and Fitters"
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.color.rgb = DARK_GRAY
p.alignment = PP_ALIGN.CENTER

# Left column: The Problem
left_title = slide17.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.5), Inches(0.5))
tf = left_title.text_frame
tf.text = "The Problem We Solved"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = NAVY

left_content = slide17.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(5.5), Inches(2.5))
tf = left_content.text_frame
tf.word_wrap = True
tf.text = "Technical outputs don't help golfers improve:"
p = tf.paragraphs[0]
p.font.size = Pt(14)
p.font.bold = True

bullet_points = [
    '• "Your lag angle is 16.5°" → So what?',
    '• "R² = 0.849" → What does that mean for my slice?',
    '• "X-Factor = 34°" → Should I stretch? Rotate more?',
    '',
    'The Translation Layer converts ML outputs into',
    'actionable golf intelligence.'
]

for point in bullet_points:
    p = tf.add_paragraph()
    p.text = point
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    if point.startswith('•'):
        p.level = 1
    elif point.startswith('The Translation'):
        p.font.bold = True
        p.font.color.rgb = NAVY

# Right column: Three Dashboards
right_title = slide17.shapes.add_textbox(Inches(6.5), Inches(1.8), Inches(6), Inches(0.5))
tf = right_title.text_frame
tf.text = "Three User-Facing Dashboards"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = NAVY

# Dashboard 1: Golfer
db1_title = slide17.shapes.add_textbox(Inches(6.5), Inches(2.3), Inches(6), Inches(0.4))
tf = db1_title.text_frame
tf.text = "🏌️ Golfer Report"
p = tf.paragraphs[0]
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = GOLD

db1_content = slide17.shapes.add_textbox(Inches(6.5), Inches(2.7), Inches(6), Inches(0.6))
tf = db1_content.text_frame
tf.word_wrap = True
tf.text = '"You\'re losing 22 yards. Do the Pump Drill daily."'
p = tf.paragraphs[0]
p.font.size = Pt(11)
p.font.italic = True
p.font.color.rgb = DARK_GRAY

# Dashboard 2: Coach
db2_title = slide17.shapes.add_textbox(Inches(6.5), Inches(3.4), Inches(6), Inches(0.4))
tf = db2_title.text_frame
tf.text = "👨‍🏫 Coach Dashboard"
p = tf.paragraphs[0]
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = GOLD

db2_content = slide17.shapes.add_textbox(Inches(6.5), Inches(3.8), Inches(6), Inches(0.6))
tf = db2_content.text_frame
tf.word_wrap = True
tf.text = '"Student improved 5° this month. Add speed training."'
p = tf.paragraphs[0]
p.font.size = Pt(11)
p.font.italic = True
p.font.color.rgb = DARK_GRAY

# Dashboard 3: Fitter
db3_title = slide17.shapes.add_textbox(Inches(6.5), Inches(4.5), Inches(6), Inches(0.4))
tf = db3_title.text_frame
tf.text = "⛳ Equipment Fitter"
p = tf.paragraphs[0]
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = GOLD

db3_content = slide17.shapes.add_textbox(Inches(6.5), Inches(4.9), Inches(6), Inches(0.6))
tf = db3_content.text_frame
tf.word_wrap = True
tf.text = '"Recommend X-stiff shaft for late release pattern."'
p = tf.paragraphs[0]
p.font.size = Pt(11)
p.font.italic = True
p.font.color.rgb = DARK_GRAY

# Bottom: Impact comparison
impact_title = slide17.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(0.5))
tf = impact_title.text_frame
tf.text = "The Impact: From Research Project to Product Golfers Use"
p = tf.paragraphs[0]
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = NAVY
p.alignment = PP_ALIGN.CENTER

# Comparison table
left_col = slide17.shapes.add_textbox(Inches(1), Inches(6.0), Inches(5), Inches(1.2))
tf = left_col.text_frame
tf.text = "❌ WITHOUT Translation Layer"
p = tf.paragraphs[0]
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(180, 0, 0)

points = ['• "Your R² is 0.849"', '• Technical reports', '• Confusing metrics']
for pt in points:
    p = tf.add_paragraph()
    p.text = pt
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY

right_col = slide17.shapes.add_textbox(Inches(7), Inches(6.0), Inches(5.5), Inches(1.2))
tf = right_col.text_frame
tf.text = "✅ WITH Translation Layer"
p = tf.paragraphs[0]
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 128, 0)

points = ['• "You\'re losing 22 yards. Here\'s the drill."', '• 7-day practice plans', '• Plain English + expected outcomes']
for pt in points:
    p = tf.add_paragraph()
    p.text = pt
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY

print("✅ Slide 17 created: Translation Layer")

# ============================================================
# SLIDE 18: Complete Platform
# ============================================================
print("Creating Slide 18: Complete Platform...")

slide18 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title_box = slide18.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
tf = title_box.text_frame
tf.text = "The Complete GolfBioMetrics Platform"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = NAVY
p.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide18.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.4))
tf = subtitle_box.text_frame
tf.text = "Technical Foundation + Translation Layer = Product-Market Fit"
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.color.rgb = GOLD
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Layer 1: Technical Excellence
layer1_title = slide18.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(6), Inches(0.5))
tf = layer1_title.text_frame
tf.text = "Layer 1: Technical Excellence ✅"
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY

layer1_content = slide18.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(6), Inches(1.5))
tf = layer1_content.text_frame
tf.word_wrap = True
points = [
    '• 58 production features across 4 domains',
    '• 5 ML models with R² 0.85–0.96',
    '• 0 NaN, 0 infinity, 0 duplicates',
    '• SHAP explainability for every prediction'
]
tf.text = points[0]
p = tf.paragraphs[0]
p.font.size = Pt(12)
for pt in points[1:]:
    p = tf.add_paragraph()
    p.text = pt
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY

# Layer 2: Translation Layer
layer2_title = slide18.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(6), Inches(0.5))
tf = layer2_title.text_frame
tf.text = "Layer 2: Translation Layer ✅"
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY

layer2_content = slide18.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(6), Inches(1.5))
tf = layer2_content.text_frame
tf.word_wrap = True
points = [
    '• Golfer dashboard with drill database',
    '• Coach dashboard with progress tracking',
    '• Equipment fitter with physics-based recs',
    '• Unified API: for_golfer() / for_coach() / for_fitter()'
]
tf.text = points[0]
p = tf.paragraphs[0]
p.font.size = Pt(12)
for pt in points[1:]:
    p = tf.add_paragraph()
    p.text = pt
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY

# Revenue Streams Section
revenue_title = slide18.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12), Inches(0.5))
tf = revenue_title.text_frame
tf.text = "Revenue Streams Enabled by Translation Layer"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = NAVY
p.alignment = PP_ALIGN.CENTER

# Revenue table
revenue_data = [
    ("Golfers", "Personal Swing Coach App", "$29.99/mo", "Drill recommendations, practice plans"),
    ("Coaches", "Pro Dashboard", "$99/mo", "Student management, data-driven insights"),
    ("Equipment", "Fitting API", "$0.50/fitting", "Physics-based recommendations"),
]

table_left = Inches(1)
table_top = Inches(4.2)
table_width = Inches(11)
table_height = Inches(1.5)

table = slide18.shapes.add_table(4, 4, table_left, table_top, table_width, table_height).table

# Headers
headers = ['User Type', 'Product', 'Price', 'Translation Layer Role']
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(11)
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Data rows
for row_idx, row_data in enumerate(revenue_data, 1):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = cell_text
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        if col_idx == 0:  # User Type column
            cell.text_frame.paragraphs[0].font.bold = True

# Quote Section
quote_box = slide18.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12), Inches(1.0))
tf = quote_box.text_frame
tf.word_wrap = True
tf.text = '"We\'ve built the most accurate golf ML system available. But accuracy alone doesn\'t create value. The Translation Layer turns \'R² = 0.849\' into \'practice this drill for 10 minutes.\' This is how we achieve product-market fit and build a $24M ARR business."'
p = tf.paragraphs[0]
p.font.size = Pt(12)
p.font.italic = True
p.font.color.rgb = DARK_GRAY
p.alignment = PP_ALIGN.CENTER

# Status
status_box = slide18.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
tf = status_box.text_frame
tf.text = "Status: Production Ready — Awaiting DSG Deployment Approval"
p = tf.paragraphs[0]
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 128, 0)
p.alignment = PP_ALIGN.CENTER

print("✅ Slide 18 created: Complete Platform")

# ============================================================
# SAVE PRESENTATION
# ============================================================
print()
print("Saving updated presentation...")

prs.save(pptx_path)

new_slide_count = len(prs.slides)
print(f"✅ Saved: {pptx_path}")
print(f"📊 New slide count: {new_slide_count} (added {new_slide_count - current_slide_count} slides)")
print()

# Update presentation README
readme_content = """# GolfBioMetrics - Master Presentation Package

**Unified presentation for Data Sports Group (DSG)**  
**Date:** June 18, 2026  
**Status:** Production Ready

---

## Master Presentation Files

### 1. Primary PowerPoint (18 Slides) - UPDATED
**`GolfBioMetrics_DSG_Master_Presentation.pptx`**
- Professional navy/gold theme
- 18 comprehensive slides covering all aspects - NOW WITH 2 NEW SLIDES
- Suitable for all audiences (executives, technical, business)

**Slide Structure:**
1. Title slide
2. Executive Summary (5-model results)
3. Problem & Market Gap
4. 3-Layer Architecture
5. Dataset Overview (500 swings, 626K frames)
6. Feature Engineering (58 features, 4 domains)
7. The 7 Core Biomechanics Metrics
8. Demographics & Environmental Differentiators
9. Models 1 & 2 (Linear Regression + Decision Tree)
10. Models 3 & 4 (Random Forest + XGBoost/SHAP)
11. Model 5 & Comparison Table (SVM + Full Results)
12. SHAP Explainability (3 Golfer Profiles)
13. Statistical Validation
14. Business Model ($24M ARR)
15. Competitive Moat & 90-Day Roadmap
16. Technical Appendix
17. **NEW: Translation Layer** (ML to Human Communication)
18. **NEW: Complete Platform** (Product-Market Fit Achieved)

### 2. Full Markdown Narrative (40 KB) ⬅️ UPDATED
**`GolfBioMetrics_DSG_Master_Presentation.md`**
- Complete speaker notes and talking points
- Detailed explanations for every slide (now 18 slides)
- Q&A preparation included
- Ready for web presentation tools (Marp, Reveal.js)

---

## Supporting Visual Assets

Located in `outputs/figures/`:

| File | Description | Use Case |
|------|-------------|----------|
| `premium_championship_infographic.png` | 6-chart dashboard (architecture, features, performance, revenue, timeline) | Full-page slide insert |
| `premium_executive_dashboard.png` | Executive metrics (system scale, data quality, accuracy, revenue) | Board meeting visuals |
| `time_series_statistical_analysis.png` | Advanced statistics (velocity, jerk, FFT, lag features) | Technical deep-dive |
| `age_experience_analysis.png` | Demographics impact on performance | Personalization story |
| `environmental_effects_analysis.png` | Weather/course condition effects | Tournament prep angle |

---

## How to Use

### For PowerPoint Presentation:
1. **Open:** `GolfBioMetrics_DSG_Master_Presentation.pptx` (now 18 slides!)
2. **Present:** Using PowerPoint, Google Slides, or Keynote
3. **Reference:** `GolfBioMetrics_DSG_Master_Presentation.md` for speaker notes
4. **Insert:** Visual assets as full-page slides where needed

### For Web/Online:
- Convert markdown to HTML using tools like Marp or Reveal.js
- Embed visual assets as images
- Share link with stakeholders

### For Print:
- Print PPT slides (18 pages) - UPDATED
- Include visual assets as appendix
- Provide markdown narrative as leave-behind

---

## Presentation Flow Recommendation

| Time | Section | Slides | Key Message |
|------|---------|--------|-------------|
| 0-2 min | Opening | 1-2 | This is the most comprehensive golf AI system |
| 2-5 min | Problem | 3 | 65M golfers lack access to professional biomechanics |
| 5-8 min | Solution | 4-8 | 58 features, 5 models, fully validated |
| 8-12 min | Results | 9-13 | R-squared 0.85-0.96, 82-92% accuracy, SHAP explainable |
| 12-15 min | Business | 14-15 | $24M ARR potential, 90-day roadmap |
| 15-18 min | Q&A | 16-18 | Technical appendix + Translation Layer + Complete Platform |

**NEW SLIDES 17-18:** Focus on how we solve the communication problem and achieve product-market fit.

---

## File Locations

```
presentation/
├── GolfBioMetrics_DSG_Master_Presentation.pptx    <- Main deck (18 slides - UPDATED!)
├── GolfBioMetrics_DSG_Master_Presentation.md      <- Full narrative (18 slides - UPDATED!)
└── README.md                                      <- This file (UPDATED!)

outputs/figures/
├── premium_championship_infographic.png           <- Key visual asset
├── premium_executive_dashboard.png                <- Executive summary visual
└── [other analysis charts]
```

---

## Production Status

- [x] All 58 features integrated and validated
- [x] 5 ML models trained (R-squared 0.83-0.96)
- [x] Data quality confirmed (0 NaN/inf/dup)
- [x] Presentation finalized (18 slides - UPDATED June 18, 2026)
- [x] **Translation Layer implemented** (NEW!)
- [x] Visual assets created (15 figures)
- [x] Documentation complete (12+ files)

**Ready for DSG presentation!**

---

**What's New (June 18, 2026):**
- NEW Slide 17: Translation Layer - ML to Human Communication
- NEW Slide 18: Complete Platform - Product-Market Fit
- NEW 4 Python modules for Translation Layer (2,000+ lines)
- NEW Drill database with 10+ specific drills
- NEW 3 user dashboards: Golfer, Coach, Equipment Fitter

**This turns R-squared 0.849 into "practice this drill for 10 minutes."**
"""

with open('presentation/README.md', 'w') as f:
    f.write(readme_content)

print("✅ Updated: presentation/README.md")
print()
print("=" * 70)
print("SUCCESS! MASTER PRESENTATION UPDATED")
print("=" * 70)
print()
print("Summary:")
print(f"  📊 Old slide count: {current_slide_count}")
print(f"  📊 New slide count: {new_slide_count}")
print(f"  ➕ Added: {new_slide_count - current_slide_count} new slides")
print()
print("New Slides:")
print("  Slide 17: The Translation Layer")
print("           └─ From ML Metrics to Human Understanding")
print("  Slide 18: The Complete Platform")
print("           └─ Technical Foundation + Translation Layer = Product-Market Fit")
print()
print("Files Updated:")
print("  ✅ presentation/GolfBioMetrics_DSG_Master_Presentation.pptx")
print("  ✅ presentation/README.md")
print()
print("Next: Push to GitHub to sync changes")
print("=" * 70)
